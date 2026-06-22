from pptx import Presentation
from pptx.util import Pt
from typing import List
from .base_parser import BaseParser, ParsedChunk


class PPTXParser(BaseParser):
    def parse(self, file_path: str) -> List[ParsedChunk]:
        chunks = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_chunks = self._extract_slide(slide, file_path, slide_num)
            chunks.extend(slide_chunks)
        return chunks

    def _extract_slide(self, slide, file_path: str,
                       slide_num: int) -> List[ParsedChunk]:
        chunks = []
        title = ""
        body_texts = []
        notes_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.shape_type == 13:  # picture
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Detect title by placeholder type
                if hasattr(shape, "placeholder_format") and \
                   shape.placeholder_format and \
                   shape.placeholder_format.idx == 0:
                    title = text
                else:
                    body_texts.append(text)

            # Extract tables inside slides
            if shape.has_table:
                table = shape.table
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in table.rows
                ]
                if len(rows) < 2:
                    continue
                headers = rows[0]
                for row_idx, row in enumerate(rows[1:], start=1):
                    row_text = " | ".join(
                        f"{headers[i]}={val}"
                        for i, val in enumerate(row) if val
                    )
                    if not row_text.strip():
                        continue
                    meta = self._base_metadata(file_path, "pptx")
                    meta.update({
                        "slide_number": slide_num,
                        "slide_title": title,
                        "row_index": row_idx,
                        "headers": headers,
                        "chunk_type": "table_row",
                    })
                    chunks.append(ParsedChunk(text=row_text, metadata=meta))

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # Combine title + body as one prose chunk per slide
        full_text_parts = []
        if title:
            full_text_parts.append(f"Slide Title: {title}")
        if body_texts:
            full_text_parts.append("\n".join(body_texts))
        if notes_text:
            full_text_parts.append(f"Speaker Notes: {notes_text}")

        if full_text_parts:
            meta = self._base_metadata(file_path, "pptx")
            meta.update({
                "slide_number": slide_num,
                "slide_title": title,
                "chunk_type": "slide",
            })
            chunks.append(ParsedChunk(
                text="\n".join(full_text_parts),
                metadata=meta
            ))

        return chunks